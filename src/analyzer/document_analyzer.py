import os
import mimetypes
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image

class DocumentAnalyzer:
    def __init__(self):
        mimetypes.init()
        # Add additional MIME types that might not be in the default database
        mimetypes.add_type('application/vnd.openxmlformats-officedocument.wordprocessingml.document', '.docx')
        mimetypes.add_type('application/vnd.openxmlformats-officedocument.presentationml.presentation', '.pptx')
        mimetypes.add_type('application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', '.xlsx')

    def analyze_file(self, file_path):
        """
        Analyze different types of documents and return relevant information
        """
        try:
            mime_type, _ = mimetypes.guess_type(file_path)
            if mime_type is None:
                # Default to binary if type cannot be determined
                mime_type = 'application/octet-stream'
            
            file_size = os.path.getsize(file_path)
            
            result = {
                'file_type': mime_type,
                'file_size': file_size,
                'analysis': {}
            }

            if file_path.lower().endswith('.docx'):
                result['analysis'] = self._analyze_word(file_path)
            elif file_path.lower().endswith('.pptx'):
                result['analysis'] = self._analyze_powerpoint(file_path)
            elif file_path.lower().endswith('.xlsx'):
                result['analysis'] = self._analyze_excel(file_path)
            elif mime_type and mime_type.startswith('image/'):
                result['analysis'] = self._analyze_image(file_path)
            elif mime_type and mime_type.startswith('text/'):
                result['analysis'] = self._analyze_text(file_path)

            return result
        except Exception as e:
            return {'error': f"Error analyzing document: {str(e)}"}

    def _analyze_word(self, file_path):
        """
        Analyze Word document
        """
        try:
            doc = Document(file_path)
            return {
                'paragraphs': len(doc.paragraphs),
                'sections': len(doc.sections),
                'tables': len(doc.tables),
                'has_macros': self._check_macros(file_path)
            }
        except:
            return {}

    def _analyze_powerpoint(self, file_path):
        """
        Analyze PowerPoint presentation
        """
        try:
            prs = Presentation(file_path)
            return {
                'slides': len(prs.slides),
                'shapes': sum(len(slide.shapes) for slide in prs.slides),
                'has_macros': self._check_macros(file_path)
            }
        except:
            return {}

    def _analyze_excel(self, file_path):
        """
        Analyze Excel workbook
        """
        try:
            wb = load_workbook(file_path, read_only=True)
            return {
                'sheets': len(wb.sheetnames),
                'sheet_names': wb.sheetnames,
                'has_macros': self._check_macros(file_path)
            }
        except:
            return {}

    def _analyze_image(self, file_path):
        """
        Analyze image file
        """
        try:
            with Image.open(file_path) as img:
                return {
                    'format': img.format,
                    'mode': img.mode,
                    'size': img.size,
                    'info': img.info
                }
        except:
            return {}

    def _analyze_text(self, file_path):
        """
        Analyze text file
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                lines = content.split('\n')
                return {
                    'lines': len(lines),
                    'characters': len(content),
                    'words': len(content.split())
                }
        except:
            return {}

    def _check_macros(self, file_path):
        """
        Check if Office document contains macros
        """
        try:
            # Simple check for macro presence by looking for .bin files in zip
            import zipfile
            with zipfile.ZipFile(file_path) as zf:
                return any(name.endswith('.bin') for name in zf.namelist())
        except:
            return None 